from __future__ import annotations

import json
import sys
import zipfile
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "backend"))
import app as app_module  # noqa: E402


class FakeCollection:
    def __init__(self) -> None:
        self.deleted: list[dict | None] = []
        self.added: list[dict] = []

    def delete(self, where=None) -> None:
        self.deleted.append(where)

    def add(self, ids, documents, metadatas) -> None:
        self.added.append({
            "ids": ids,
            "documents": documents,
            "metadatas": metadatas,
        })


def _write_minimal_docx(path: Path) -> None:
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            (
                '<?xml version="1.0"?>'
                '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
                '<Default Extension="xml" ContentType="application/xml"/>'
                '<Override PartName="/word/document.xml" '
                'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
                "</Types>"
            ),
        )
        archive.writestr(
            "_rels/.rels",
            (
                '<?xml version="1.0"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" '
                'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
                'Target="word/document.xml"/>'
                "</Relationships>"
            ),
        )
        archive.writestr(
            "word/document.xml",
            (
                '<?xml version="1.0"?>'
                '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
                "<w:body>"
                "<w:p><w:r><w:t>DOCX work instruction</w:t></w:r></w:p>"
                "<w:p><w:r><w:t>Torque bolts to 12 Nm.</w:t></w:r></w:p>"
                "</w:body></w:document>"
            ),
        )


def _write_minimal_pdf(path: Path) -> None:
    path.write_bytes(
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Count 1/Kids[3 0 R]>>endobj\n"
        b"3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 300 144]/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj\n"
        b"4 0 obj<</Length 74>>stream\n"
        b"BT /F1 18 Tf 36 96 Td (PDF safety inspection procedure) Tj ET\n"
        b"endstream endobj\n"
        b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"xref\n0 6\n0000000000 65535 f \n"
        b"0000000009 00000 n \n0000000058 00000 n \n0000000115 00000 n \n"
        b"0000000234 00000 n \n0000000358 00000 n \n"
        b"trailer<</Size 6/Root 1 0 R>>\nstartxref\n426\n%%EOF\n"
    )


def _write_workbook(path: Path) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Calibration"
    sheet.append(["Asset", "Status", "Note"])
    sheet.append(["Gauge A", "Active", "Calibrated"])
    workbook.save(path)


def _write_presentation(path: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "PPTX maintenance briefing"
    text_box = slide.shapes.add_textbox(914400, 1371600, 7315200, 914400)
    text_box.text_frame.text = "Inspect guards before startup."
    presentation.save(path)


@pytest.fixture()
def supported_samples(tmp_path: Path) -> dict[str, tuple[Path, str]]:
    samples: dict[str, tuple[Path, str]] = {}

    text_samples = {
        ".txt": ("sample.txt", "TXT lockout procedure"),
        ".md": ("sample.md", "# MD startup guide"),
        ".log": ("sample.log", "LOG error E-100 resolved"),
        ".json": ("sample.json", json.dumps({"title": "JSON calibration procedure"})),
        ".xml": ("sample.xml", "<root><title>XML inspection checklist</title></root>"),
        ".rtf": ("sample.rtf", r"{\rtf1\ansi RTF cleaning procedure}"),
        ".csv": ("sample.csv", "part,qty,note\nBearing,4,CSV inventory row\n"),
        ".tsv": ("sample.tsv", "asset\tstatus\tnote\nPump A\tactive\tTSV calibration row\n"),
        ".html": ("sample.html", "<html><body><h1>HTML maintenance page</h1></body></html>"),
        ".htm": ("sample.htm", "<html><body><h1>HTM inspection page</h1></body></html>"),
    }
    expected_by_ext = {
        ".txt": "TXT lockout procedure",
        ".md": "MD startup guide",
        ".log": "LOG error E-100",
        ".json": "JSON calibration procedure",
        ".xml": "XML inspection checklist",
        ".rtf": "RTF cleaning procedure",
        ".csv": "CSV inventory row",
        ".tsv": "TSV calibration row",
        ".html": "HTML maintenance page",
        ".htm": "HTM inspection page",
    }

    for ext, (name, content) in text_samples.items():
        path = tmp_path / name
        path.write_text(content, encoding="utf-8")
        samples[ext] = (path, expected_by_ext[ext])

    docx_path = tmp_path / "sample.docx"
    _write_minimal_docx(docx_path)
    samples[".docx"] = (docx_path, "DOCX")

    pdf_path = tmp_path / "sample.pdf"
    _write_minimal_pdf(pdf_path)
    samples[".pdf"] = (pdf_path, "PDF")

    xlsx_path = tmp_path / "sample.xlsx"
    _write_workbook(xlsx_path)
    samples[".xlsx"] = (xlsx_path, "Gauge A")

    xlsm_path = tmp_path / "sample.xlsm"
    _write_workbook(xlsm_path)
    samples[".xlsm"] = (xlsm_path, "Gauge A")

    pptx_path = tmp_path / "sample.pptx"
    _write_presentation(pptx_path)
    samples[".pptx"] = (pptx_path, "PPTX")

    return samples


def test_all_allowed_extensions_have_sample_coverage(supported_samples):
    assert set(supported_samples) == app_module.ALLOWED_EXTENSIONS


@pytest.mark.parametrize("ext", sorted(app_module.ALLOWED_EXTENSIONS))
def test_supported_extension_loads_and_chunks(ext, supported_samples):
    path, expected = supported_samples[ext]

    docs = app_module._load_docs(path)
    chunks = app_module._prepare_chunks(path, docs)
    extracted = "\n".join(doc.page_content for doc in docs)

    assert docs
    assert chunks
    assert expected.lower() in extracted.lower()


@pytest.mark.parametrize("ext", sorted(app_module.ALLOWED_EXTENSIONS))
def test_supported_extension_ingests_into_collection(ext, supported_samples, monkeypatch):
    path, _ = supported_samples[ext]
    fake_collection = FakeCollection()
    monkeypatch.setattr(app_module, "collection", fake_collection)

    ok, message, chunk_count = app_module._ingest_file(path)

    assert ok, message
    assert chunk_count > 0
    assert fake_collection.deleted
    assert sum(len(batch["ids"]) for batch in fake_collection.added) == chunk_count
    assert fake_collection.added[0]["metadatas"][0]["source"] == str(path)
