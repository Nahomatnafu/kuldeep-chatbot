"""
Kuldeep RAG Chatbot — Flask Backend
=====================================
RAG pipeline powered by ChromaDB + LangChain + OpenAI.

Endpoints
---------
POST   /chat                         — Main chat (used by Next.js proxy)
GET    /api/documents                — List uploaded documents
POST   /api/documents/upload         — Upload + ingest a document
DELETE /api/documents/<filename>     — Delete a document + its vectors
POST   /api/clear                    — Clear conversation history for a session
GET    /api/health                   — Health check
"""

import os
import re
import csv
import json
import time
import traceback
import math
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path

from flask import Flask, request, jsonify
from flask_cors import CORS
from werkzeug.utils import secure_filename
from dotenv import load_dotenv

import chromadb
from chromadb.utils.embedding_functions import OpenAIEmbeddingFunction

from langchain_openai import ChatOpenAI
from langchain.prompts import PromptTemplate
from langchain.memory import ConversationBufferMemory
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

load_dotenv()

# ── Configuration ────────────────────────────────────────────────────────────
KNOWLEDGE_BASE_DIR = Path("knowledge_base")
CHROMA_DB_DIR      = "chroma_db"
DOCUMENTS_JSON     = Path("knowledge_base/documents.json")
MODEL_NAME         = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
EMBEDDING_MODEL    = os.getenv("OPENAI_EMBEDDING_MODEL", "text-embedding-ada-002")
NUM_CHUNKS         = 12         # k=12 for broader recall on dense technical documents
CHUNK_SIZE         = 1000
CHUNK_OVERLAP      = 200
VECTOR_POOL_FACTOR = 4          # retrieve wider, then rerank with lexical signals
KEYWORD_POOL_SIZE  = 40
HYBRID_ALPHA       = 0.72       # vector weight; lexical/exact-match fills the rest
MAX_UPLOAD_MB      = 32
ALLOWED_EXTENSIONS = {
    ".pdf", ".txt", ".md", ".log", ".json", ".xml", ".docx", ".rtf",
    ".csv", ".tsv", ".xlsx", ".xlsm", ".pptx", ".html", ".htm",
}
BROAD_INTENT_PHRASES = {
    "all documents", "all the documents", "all docs", "all the docs",
    "across all", "across documents", "across the documents", "from all",
    "compare", "consolidate", "summarize all", "all relevant",
    "every document", "every doc", "from every",
}
SESSION_TTL_SECONDS      = 2 * 60 * 60   # evict sessions idle for 2 hours
MAX_MULTI_DOC_CHUNKS     = 40            # hard cap on total chunks sent to LLM (≈10K tokens)

KNOWLEDGE_BASE_DIR.mkdir(exist_ok=True)


def _collection_name_for_embedding(model_name: str) -> str:
    """Use one Chroma collection per embedding model to avoid dimension conflicts."""
    if model_name == "text-embedding-ada-002":
        return "documents"
    safe = re.sub(r"[^a-zA-Z0-9_]+", "_", model_name).strip("_").lower()
    return f"documents_{safe}"


COLLECTION_NAME = _collection_name_for_embedding(EMBEDDING_MODEL)

# ── Flask app ─────────────────────────────────────────────────────────────────
app = Flask(__name__)
CORS(app)
app.config["MAX_CONTENT_LENGTH"] = MAX_UPLOAD_MB * 1024 * 1024

# ── Global state ─────────────────────────────────────────────────────────────
collection: chromadb.Collection | None = None
llm:        ChatOpenAI | None          = None
_guard_llm: ChatOpenAI | None          = None   # cheap off-topic classifier
conversation_sessions: dict            = {}      # session_id → { memory, last_accessed }
_keyword_corpus_cache: dict            = {}

# ── Prompt ───────────────────────────────────────────────────────────────────
_QA_TEMPLATE = """You are a careful research assistant. \
Your ONLY job is to extract and summarize information from the provided context.

STRICT GROUNDING RULES (DO NOT VIOLATE):
1. Answer ONLY from the retrieved context below — do NOT add facts not present in it.
2. Every claim in your answer MUST be traceable to a specific chunk below.
3. If the context is insufficient or doesn't contain the answer → say ONLY: \
"The uploaded documents do not contain information about this."
4. Do NOT infer, extrapolate, or fill gaps with general knowledge.
5. For lists (models, methods, steps, names) → ONLY list what is explicitly named \
in the context.
6. If the question is broad but context is limited → acknowledge the limitation: \
"Based on the provided context, the following are mentioned: ..."
7. Your answer should read like a direct summary of the chunks below, \
NOT like an essay.

FORMATTING RULES (always apply):
- Use **bold** to highlight key terms, names, model numbers, or critical warnings.
- Use bullet points or numbered lists whenever presenting multiple items, steps, or options.
- Use numbered steps for sequential procedures (e.g. startup, shutdown, inspection).
- Use short paragraphs — avoid walls of text.
- If the answer has clearly distinct sections, add a short **bold** heading for each.
- Keep formatting clean and readable — do not over-format simple one-line answers.

Retrieved Context:
{context}

Question: {question}

Answer (extract/summarize ONLY from context above, applying formatting rules):

After your answer, on a separate line write exactly:
SOURCES_USED: followed by comma-separated chunk numbers you drew from (e.g. SOURCES_USED: 1, 3, 5).
If you used none of the chunks, write: SOURCES_USED: none"""

_CONDENSE_TEMPLATE = """Given the following conversation and a follow up question, \
rephrase the follow up question to be a standalone question ONLY IF truly needed.

Rules:
- If the follow-up is about a NEW topic or entity not mentioned before, \
return it EXACTLY as written — do NOT inject entities or topics from prior turns.
- Only rephrase if the follow-up uses pronouns (it, they, this, that) or \
clearly refers back to something already discussed \
(e.g. "tell me more", "what about the second one", "why is that?").
- If the chat history is empty, return the question as-is.

Chat History:
{chat_history}

Follow Up Question: {question}

Standalone question:"""

QA_PROMPT       = PromptTemplate(template=_QA_TEMPLATE,       input_variables=["context", "question"])
CONDENSE_PROMPT = PromptTemplate(template=_CONDENSE_TEMPLATE, input_variables=["chat_history", "question"])

_MULTI_DOC_QA_TEMPLATE = """\
You are a careful research assistant synthesizing information from multiple documents.

RULES:
1. Answer ONLY from the source chunks provided below — do NOT add external knowledge.
2. Cite which document each point comes from using [Source: filename].
3. If documents differ or conflict on a point, say so explicitly.
4. If the context is insufficient → say "The uploaded documents do not contain enough information about this."

FORMATTING RULES (always apply):
- Use **bold** to highlight key terms, names, model numbers, or critical warnings.
- Use bullet points or numbered lists whenever presenting multiple items, steps, or options.
- Use numbered steps for sequential procedures (e.g. startup, shutdown, inspection).
- Group information by document or topic using short **bold** headings where it aids clarity.
- Use short paragraphs — avoid walls of text.
- Keep formatting clean and readable — do not over-format simple one-line answers.

Retrieved Context (grouped by source):
{context}

Question: {question}

Answer (synthesize from context above, citing sources, applying formatting rules):

After your answer, on a separate line write exactly:
SOURCES_USED: followed by comma-separated chunk numbers you drew from (e.g. SOURCES_USED: 2, 4).
If you used none of the chunks, write: SOURCES_USED: none"""

# ── Off-topic guard ───────────────────────────────────────────────────────────
_GUARD_PROMPT = (
    "You are a classifier for a workplace document assistant. "
    "Answer with exactly one word: YES or NO.\n\n"
    "Answer YES only if the message is clearly one of these three things:\n"
    "1. A question about the user's own identity or personal life "
    "(e.g. 'What is my name?', 'How old am I?')\n"
    "2. Pure social small talk with no factual question "
    "(e.g. 'How are you?', 'Good morning')\n"
    "3. A request for jokes, stories, poems, songs, or other entertainment "
    "(e.g. 'Tell me a joke', 'Write me a poem')\n\n"
    "Answer NO for everything else, including any factual, technical, or "
    "procedural question — even if it seems unrelated to the documents. "
    "The document system will handle those appropriately on its own.\n\n"
    "When in doubt, answer NO.\n\n"
    "User message: \"{msg}\"\n\n"
    "Answer (YES or NO):"
)

def _is_off_topic(question: str) -> bool:
    """
    Make a cheap, context-free LLM call to classify whether the question is
    clearly personal / off-topic BEFORE the RAG chain retrieves any documents.
    Running this *before* retrieval prevents document context from biasing the answer.
    Returns True  → block the question and return the canned off-topic reply.
    Returns False → let the full RAG chain handle it normally.
    """
    # Fast deterministic check for obvious entertainment requests (whole-word match)
    _lower = question.lower()
    _ENTERTAINMENT_WORDS = ("joke", "jokes", "story", "stories", "poem", "poems",
                            "riddle", "riddles", "sing", "song", "limerick")
    if any(re.search(r'\b' + w + r'\b', _lower) for w in _ENTERTAINMENT_WORDS):
        return True

    global _guard_llm
    try:
        if _guard_llm is None:
            _guard_llm = ChatOpenAI(
                model_name="gpt-4o-mini",
                temperature=0,
                max_tokens=3,      # only need "YES" or "NO"
            )
        safe_msg = question.replace('"', "'")
        result   = _guard_llm.invoke(_GUARD_PROMPT.format(msg=safe_msg))
        raw      = result.content.strip()
        return raw.upper().startswith("YES")
    except Exception:
        return False   # fail open — let the chain handle edge cases


# ── Helper: document registry ─────────────────────────────────────────────────
def _load_doc_registry() -> dict:
    """Load the JSON registry of ingested documents."""
    if DOCUMENTS_JSON.exists():
        return json.loads(DOCUMENTS_JSON.read_text())
    return {}


def _save_doc_registry(registry: dict) -> None:
    DOCUMENTS_JSON.write_text(json.dumps(registry, indent=2))


def _clean_text(text: str) -> str:
    """Normalize whitespace while preserving enough structure for headings."""
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_revision_status(text: str, filename: str) -> tuple[str, str, str]:
    """Best-effort SOP/manual metadata from filenames and first-page text."""
    sample = f"{filename}\n{text[:4000]}"
    revision = ""
    effective_date = ""
    status = "active"

    rev_match = re.search(r"\b(?:rev(?:ision)?\.?|version|ver\.?)\s*[:#-]?\s*([A-Z0-9.\-]+)", sample, re.I)
    if rev_match:
        revision = rev_match.group(1).strip(" .:-")

    date_match = re.search(
        r"\b(?:effective date|eff\.? date|issued|approved date)\s*[:#-]?\s*([A-Za-z0-9,/\-. ]{6,30})",
        sample,
        re.I,
    )
    if date_match:
        effective_date = date_match.group(1).strip(" .")

    if re.search(r"\b(obsolete|superseded|retired|archived|do not use)\b", sample, re.I):
        status = "obsolete"

    return revision, effective_date, status


def _derive_doc_metadata(filepath: Path, docs: list[Document]) -> dict:
    """Attach business-friendly metadata used by filters, citations, and reranking."""
    combined = "\n".join(d.page_content for d in docs[:3])
    revision, effective_date, status = _extract_revision_status(combined, filepath.name)
    stem_words = Path(filepath.name).stem.lower().replace("-", " ").replace("_", " ")
    doc_type = "document"
    if "sop" in stem_words or "procedure" in stem_words:
        doc_type = "sop"
    elif "manual" in stem_words or "operation" in stem_words or "guide" in stem_words:
        doc_type = "manual"
    elif "checklist" in stem_words:
        doc_type = "checklist"
    elif "audit" in stem_words:
        doc_type = "audit"

    return {
        "filename": filepath.name,
        "document_title": Path(filepath.name).stem.replace("_", " ").replace("-", " "),
        "doc_type": doc_type,
        "revision": revision,
        "effective_date": effective_date,
        "status": status,
        "embedding_model": EMBEDDING_MODEL,
    }


_SECTION_HEADING_RE = re.compile(
    r"^\s*(?:"
    r"(?:\d+(?:\.\d+){0,4})\s+[\w][\w ,/&()\-:]{3,}"
    r"|(?:[A-Z][A-Z0-9 /&()\-:]{5,})"
    r"|(?:scope|purpose|procedure|responsibilit(?:y|ies)|safety|warning|caution|steps?|materials?|equipment|references?)\b.*"
    r")\s*$",
    re.I,
)


def _split_docs_by_section(docs: list[Document], base_metadata: dict) -> list[Document]:
    """
    Create section-aware parent documents before character chunking.
    This keeps SOP headings, warnings, and procedure boundaries visible to retrieval.
    """
    section_docs: list[Document] = []
    current_heading = ""

    for doc in docs:
        text = _clean_text(doc.page_content)
        if not text:
            continue
        page = doc.metadata.get("page", 0)
        lines = text.splitlines()
        buffer: list[str] = []

        def flush() -> None:
            if not buffer:
                return
            content = _clean_text("\n".join(buffer))
            if content:
                section_docs.append(Document(
                    page_content=content,
                    metadata={
                        **base_metadata,
                        **{k: v for k, v in doc.metadata.items() if isinstance(v, (str, int, float, bool))},
                        "page": page,
                        "section": current_heading,
                    },
                ))
            buffer.clear()

        for raw_line in lines:
            line = raw_line.strip()
            if _SECTION_HEADING_RE.match(line) and len(line) <= 120:
                flush()
                current_heading = line
                buffer.append(line)
            else:
                buffer.append(raw_line)
        flush()

    if section_docs:
        merged: list[Document] = []
        carry: Document | None = None
        for section_doc in section_docs:
            content = section_doc.page_content.strip()
            if carry is not None:
                same_page = carry.metadata.get("page") == section_doc.metadata.get("page")
                if same_page:
                    carry.page_content = _clean_text(carry.page_content + "\n" + content)
                    if not carry.metadata.get("section") and section_doc.metadata.get("section"):
                        carry.metadata["section"] = section_doc.metadata.get("section", "")
                    if len(carry.page_content) < 220:
                        continue
                    merged.append(carry)
                    carry = None
                    continue
                merged.append(carry)
                carry = None

            if len(content) < 120:
                carry = section_doc
            else:
                merged.append(section_doc)
        if carry is not None:
            if merged and merged[-1].metadata.get("page") == carry.metadata.get("page"):
                merged[-1].page_content = _clean_text(merged[-1].page_content + "\n" + carry.page_content)
            else:
                merged.append(carry)
        return merged

    return [
        Document(
            page_content=_clean_text(d.page_content),
            metadata={
                **base_metadata,
                **{k: v for k, v in d.metadata.items() if isinstance(v, (str, int, float, bool))},
                "section": "",
            },
        )
        for d in docs
        if _clean_text(d.page_content)
    ]


def _prepare_chunks(filepath: Path, docs: list[Document]) -> list[Document]:
    base_metadata = _derive_doc_metadata(filepath, docs)
    section_docs = _split_docs_by_section(docs, base_metadata)
    splitter = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    chunks = splitter.split_documents(section_docs)
    for i, chunk in enumerate(chunks):
        chunk.metadata["chunk_index"] = i
        chunk.metadata["chunk_id"] = f"{filepath.name}__chunk_{i}"
    return chunks


# ── Helper: initialise / reload vector store ─────────────────────────────────
def _init_store() -> tuple[bool, str]:
    """Initialise (or reload) the ChromaDB collection and LLM. Returns (ok, message)."""
    global collection, llm

    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        return False, "OPENAI_API_KEY not set in environment."

    try:
        ef         = OpenAIEmbeddingFunction(api_key=api_key, model_name=EMBEDDING_MODEL)
        client     = chromadb.PersistentClient(path=CHROMA_DB_DIR)
        collection = client.get_or_create_collection(
            name=COLLECTION_NAME,
            embedding_function=ef,
            metadata={"embedding_model": EMBEDDING_MODEL},
        )
        llm        = ChatOpenAI(model_name=MODEL_NAME, temperature=0)
        return True, f"Vector store loaded ({COLLECTION_NAME})."
    except Exception as exc:
        return False, f"Failed to load vector store: {exc}"


# ── Helper: load documents from any supported file type ─────────────────────
def _format_pdf_table(table: list[list]) -> str:
    """Convert an extracted PDF table to compact markdown-like text."""
    cleaned_rows = []
    for row in table or []:
        cleaned = [str(cell or "").strip().replace("\n", " ") for cell in row]
        if any(cleaned):
            cleaned_rows.append(cleaned)
    if not cleaned_rows:
        return ""
    width = max(len(row) for row in cleaned_rows)
    normalized = [row + [""] * (width - len(row)) for row in cleaned_rows]
    return "\n".join("| " + " | ".join(row) + " |" for row in normalized)


def _load_pdf_docs(filepath: Path) -> list[Document]:
    """
    Prefer pdfplumber when available because it preserves tables better than
    plain PDF text extraction. Falls back to PyPDFLoader for compatibility.
    """
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        return PyPDFLoader(str(filepath)).load()

    docs: list[Document] = []
    with pdfplumber.open(str(filepath)) as pdf:
        for page_index, page in enumerate(pdf.pages):
            parts = []
            text = page.extract_text(x_tolerance=1, y_tolerance=3) or ""
            if text.strip():
                parts.append(text)
            try:
                tables = page.extract_tables() or []
            except Exception:
                tables = []
            for table_index, table in enumerate(tables, 1):
                table_text = _format_pdf_table(table)
                if table_text:
                    parts.append(f"[TABLE {table_index}]\n{table_text}")
            content = _clean_text("\n\n".join(parts))
            if content:
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": str(filepath),
                        "page": page_index,
                        "total_pages": len(pdf.pages),
                        "parser": "pdfplumber",
                    },
                ))

    return docs or PyPDFLoader(str(filepath)).load()


def _read_text_file(filepath: Path) -> str:
    """Read common text-like files with a small encoding fallback chain."""
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return filepath.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return filepath.read_text(encoding="utf-8", errors="replace")


def _load_table_docs(filepath: Path, delimiter: str) -> list[Document]:
    """Load CSV/TSV rows, supporting both headered and headerless files."""
    rows: list[str] = []
    text = _read_text_file(filepath)
    sample = text[:2048]
    try:
        has_header = csv.Sniffer().has_header(sample)
    except csv.Error:
        has_header = True

    reader = csv.reader(text.splitlines(), delimiter=delimiter)
    raw_rows = [row for row in reader if any(cell.strip() for cell in row)]
    if not raw_rows:
        return []

    if has_header and len(raw_rows) > 1:
        headers = [h.strip() or f"column_{i + 1}" for i, h in enumerate(raw_rows[0])]
        for row_index, row in enumerate(raw_rows[1:], 1):
            cells = row + [""] * max(0, len(headers) - len(row))
            values = [
                f"{header}: {value.strip()}"
                for header, value in zip(headers, cells)
                if value.strip()
            ]
            if values:
                rows.append(f"Row {row_index}: " + ", ".join(values))
    else:
        for row_index, row in enumerate(raw_rows, 1):
            values = [
                f"column_{col_index}: {value.strip()}"
                for col_index, value in enumerate(row, 1)
                if value.strip()
            ]
            if values:
                rows.append(f"Row {row_index}: " + ", ".join(values))

    return [Document(page_content="\n".join(rows), metadata={"source": str(filepath)})] if rows else []


def _load_workbook_docs(filepath: Path) -> list[Document]:
    """Extract visible worksheet values from modern Excel workbooks."""
    from openpyxl import load_workbook

    workbook = load_workbook(filename=str(filepath), read_only=True, data_only=True)
    docs: list[Document] = []
    try:
        for sheet_index, sheet in enumerate(workbook.worksheets, 1):
            rows: list[str] = []
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), 1):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    rows.append(f"Row {row_index}: " + " | ".join(values))
            content = _clean_text("\n".join(rows))
            if content:
                docs.append(Document(
                    page_content=f"Sheet: {sheet.title}\n{content}",
                    metadata={
                        "source": str(filepath),
                        "page": sheet_index - 1,
                        "sheet": sheet.title,
                        "parser": "openpyxl",
                    },
                ))
    finally:
        workbook.close()
    return docs


def _load_presentation_docs(filepath: Path) -> list[Document]:
    """Extract text from each slide in a PowerPoint deck."""
    from pptx import Presentation

    presentation = Presentation(str(filepath))
    docs: list[Document] = []
    for slide_index, slide in enumerate(presentation.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = _clean_text(shape.text)
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                table_rows: list[str] = []
                for row in shape.table.rows:
                    cells = [_clean_text(cell.text) for cell in row.cells]
                    if any(cells):
                        table_rows.append(" | ".join(cells))
                if table_rows:
                    parts.append("[TABLE]\n" + "\n".join(table_rows))
        content = _clean_text("\n\n".join(parts))
        if content:
            docs.append(Document(
                page_content=f"Slide {slide_index}\n{content}",
                metadata={
                    "source": str(filepath),
                    "page": slide_index - 1,
                    "slide": slide_index,
                    "parser": "python-pptx",
                },
            ))
    return docs


def _load_xml_docs(filepath: Path) -> list[Document]:
    """Flatten XML elements into path/value lines for retrieval."""
    import xml.etree.ElementTree as ET

    text = _read_text_file(filepath)
    try:
        root = ET.fromstring(text)
    except ET.ParseError:
        return [Document(page_content=text, metadata={"source": str(filepath)})]

    lines: list[str] = []

    def walk(element: ET.Element, path: str) -> None:
        tag = element.tag.split("}", 1)[-1]
        current_path = f"{path}/{tag}" if path else tag
        attrs = " ".join(f"@{k}: {v}" for k, v in element.attrib.items())
        value = (element.text or "").strip()
        parts = [part for part in (attrs, value) if part]
        if parts:
            lines.append(f"{current_path}: {'; '.join(parts)}")
        for child in element:
            walk(child, current_path)

    walk(root, "")
    content = "\n".join(lines) or text
    return [Document(page_content=content, metadata={"source": str(filepath)})]


def _load_docs(filepath: Path) -> list:
    """Return a list of LangChain Documents for any supported file type."""
    ext = filepath.suffix.lower()

    if ext == ".pdf":
        return _load_pdf_docs(filepath)

    if ext in (".txt", ".md", ".log"):
        return [Document(page_content=_read_text_file(filepath), metadata={"source": str(filepath)})]

    if ext == ".json":
        text = _read_text_file(filepath)
        try:
            content = json.dumps(json.loads(text), indent=2)
        except json.JSONDecodeError:
            content = text
        return [Document(page_content=content, metadata={"source": str(filepath)})]

    if ext == ".xml":
        return _load_xml_docs(filepath)

    if ext == ".docx":
        import docx2txt
        return [Document(page_content=docx2txt.process(str(filepath)), metadata={"source": str(filepath)})]

    if ext == ".rtf":
        from striprtf.striprtf import rtf_to_text
        return [Document(page_content=rtf_to_text(_read_text_file(filepath)), metadata={"source": str(filepath)})]

    if ext in (".csv", ".tsv"):
        delimiter = "\t" if ext == ".tsv" else ","
        return _load_table_docs(filepath, delimiter)

    if ext in (".xlsx", ".xlsm"):
        return _load_workbook_docs(filepath)

    if ext == ".pptx":
        return _load_presentation_docs(filepath)

    if ext in (".html", ".htm"):
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(_read_text_file(filepath), "lxml")
        for tag in soup(["script", "style"]):
            tag.decompose()
        return [Document(page_content=soup.get_text(separator="\n", strip=True), metadata={"source": str(filepath)})]

    return []


# ── Helper: ingest a single file ──────────────────────────────────────────────
def _ingest_file(filepath: Path) -> tuple[bool, str, int]:
    """Chunk and embed any supported file into ChromaDB. Returns (ok, message, chunk_count)."""
    try:
        docs = _load_docs(filepath)
        if not docs:
            return False, f"No content extracted from {filepath.name}.", 0

        chunks = _prepare_chunks(filepath, docs)

        # Remove any existing vectors for this file (handles re-uploads cleanly)
        collection.delete(where={"source": str(filepath)})

        ids       = [str(c.metadata.get("chunk_id", f"{filepath.name}__chunk_{i}")) for i, c in enumerate(chunks)]
        texts     = [c.page_content for c in chunks]
        metadatas = [{"source": str(filepath), **{k: v for k, v in c.metadata.items() if isinstance(v, (str, int, float, bool))}} for c in chunks]

        # Send in batches to stay under OpenAI's 300k token-per-request limit.
        EMBED_BATCH = 100
        for start in range(0, len(ids), EMBED_BATCH):
            end = start + EMBED_BATCH
            collection.add(ids=ids[start:end], documents=texts[start:end], metadatas=metadatas[start:end])

        _keyword_corpus_cache.clear()
        return True, f"Ingested {len(chunks)} chunks.", len(chunks)
    except Exception as exc:
        traceback.print_exc()
        return False, f"Ingestion error: {exc}", 0


# ── ChromaDB similarity search helpers ────────────────────────────────────────
def _build_numbered_context(chunks: list) -> str:
    """Format chunks as numbered sections so the LLM can cite which ones it used."""
    parts = []
    for i, c in enumerate(chunks, 1):
        src = os.path.basename(c.metadata.get("source", c.metadata.get("filename", "Unknown")))
        page = c.metadata.get("page")
        section = c.metadata.get("section", "")
        header = f"[CHUNK {i}] Source: {src}"
        if isinstance(page, int):
            header += f", Page: {page + 1}"
        if section:
            header += f", Section: {section}"
        parts.append(f"{header}\n{c.page_content}")
    return "\n\n".join(parts)


def _dedup_chunks(chunks: list) -> list:
    """Remove duplicate chunks by (source_file, page), preserving order."""
    seen = set()
    result = []
    for doc in chunks:
        key = (
            os.path.basename(doc.metadata.get("source", "")),
            doc.metadata.get("page", 0),
            doc.metadata.get("chunk_index", doc.metadata.get("chunk_id", doc.page_content[:80])),
        )
        if key not in seen:
            seen.add(key)
            result.append(doc)
    return result


def _parse_citations(raw_answer: str, chunks: list) -> tuple:
    """Parse SOURCES_USED line from answer. Returns (clean_answer, cited_chunks).
    Falls back to all chunks if the LLM omits the citation line.
    """
    match = re.search(r'\s*SOURCES_USED:\s*(.+?)\s*$', raw_answer, re.IGNORECASE | re.MULTILINE)
    if not match:
        return raw_answer.strip(), chunks

    clean_answer = raw_answer[:match.start()].strip()
    cited_str    = match.group(1).strip()

    if cited_str.lower() == "none" or not cited_str:
        return clean_answer, []

    cited_chunks = []
    for part in cited_str.split(","):
        m = re.search(r'\d+', part.strip())
        if m:
            idx = int(m.group()) - 1  # 1-based → 0-based
            if 0 <= idx < len(chunks):
                cited_chunks.append(chunks[idx])

    # Fallback: if parsing yielded nothing, return all chunks
    return clean_answer, cited_chunks if cited_chunks else chunks


_TOKEN_RE = re.compile(r"[a-zA-Z0-9][a-zA-Z0-9_.:/#-]*")
_CODE_TOKEN_RE = re.compile(
    r"(?<![A-Za-z0-9])"
    r"(?=[A-Za-z0-9_.:/#-]*[A-Za-z])"
    r"(?=[A-Za-z0-9_.:/#-]*\d)"
    r"[A-Za-z0-9][A-Za-z0-9_.:/#-]{2,}"
    r"(?![A-Za-z0-9])"
)
_LEXICAL_STOP_WORDS = {
    "the", "and", "for", "with", "what", "which", "that", "this", "from",
    "according", "document", "documents", "manual", "sop", "procedure",
    "procedures", "does", "are", "is", "to", "of", "in", "on", "a", "an",
    "all", "about", "how", "do", "i", "me", "give", "explain", "describe",
}


def _tokenize(text: str) -> list[str]:
    return [
        t.lower()
        for t in _TOKEN_RE.findall(text)
        if len(t) > 1 and t.lower() not in _LEXICAL_STOP_WORDS
    ]


def _extract_code_tokens(text: str) -> list[str]:
    """Return ordered, unique technical identifiers from a query or chunk."""
    seen: set[str] = set()
    tokens: list[str] = []
    for match in _CODE_TOKEN_RE.finditer(text):
        token = match.group(0).strip(".,;:()[]{}").lower()
        if token and token not in seen:
            seen.add(token)
            tokens.append(token)
    return tokens


def _contains_exact_identifier(haystack: str, identifier: str) -> bool:
    if not identifier:
        return False
    pattern = rf"(?<![a-z0-9]){re.escape(identifier.lower())}(?![a-z0-9])"
    return re.search(pattern, haystack.lower()) is not None


def _get_collection_docs(where: dict = None) -> list[Document]:
    """Fetch collection contents for local lexical retrieval."""
    if collection is None or collection.count() == 0:
        return []
    kwargs: dict = {"include": ["documents", "metadatas"]}
    if where:
        kwargs["where"] = where
    try:
        results = collection.get(**kwargs)
    except Exception:
        return []
    return [
        Document(page_content=t, metadata=m or {})
        for t, m in zip(results.get("documents", []), results.get("metadatas", []))
        if t
    ]


def _keyword_doc_text(doc: Document) -> str:
    return " ".join([
        doc.page_content,
        str(doc.metadata.get("filename", "")),
        str(doc.metadata.get("document_title", "")),
        str(doc.metadata.get("section", "")),
        str(doc.metadata.get("doc_type", "")),
    ])


def _keyword_corpus(where: dict = None) -> dict:
    cache_key = json.dumps(where or {}, sort_keys=True)
    count = collection.count() if collection is not None else 0
    cached = _keyword_corpus_cache.get(cache_key)
    if cached and cached.get("count") == count:
        return cached

    docs = _get_collection_docs(where=where)
    rows = []
    doc_freq: Counter = Counter()
    total_len = 0
    for doc in docs:
        haystack = _keyword_doc_text(doc)
        tokens = _tokenize(haystack)
        token_counts = Counter(tokens)
        total_len += len(tokens)
        for token in token_counts:
            doc_freq[token] += 1
        rows.append({
            "doc": doc,
            "haystack": haystack.lower(),
            "token_counts": token_counts,
            "doc_len": max(1, len(tokens)),
        })

    corpus = {
        "count": count,
        "rows": rows,
        "doc_freq": doc_freq,
        "total_docs": len(rows),
        "avg_doc_len": total_len / max(len(rows), 1),
    }
    _keyword_corpus_cache[cache_key] = corpus
    return corpus


def _keyword_score(query_tokens: list[str], row: dict, idf: dict[str, float], avg_doc_len: float) -> float:
    if not query_tokens:
        return 0.0
    haystack = row["haystack"]
    doc_tokens = row["token_counts"]
    doc_len = row["doc_len"]
    avg_doc_len = max(1.0, avg_doc_len)
    score = 0.0
    k1 = 1.5
    b = 0.75
    for token in query_tokens:
        count = doc_tokens.get(token, 0)
        if count:
            numerator = count * (k1 + 1)
            denominator = count + k1 * (1 - b + b * (doc_len / avg_doc_len))
            score += idf.get(token, 0.0) * (numerator / denominator)
        elif re.search(rf"\b{re.escape(token)}s?\b", haystack):
            score += 0.4 * idf.get(token, 0.0)
    exact_phrases = _extract_code_tokens(" ".join(query_tokens))
    for phrase in exact_phrases:
        if _contains_exact_identifier(haystack, phrase):
            score += 8.0
    return score


def _exact_identifier_search(query: str, k: int = NUM_CHUNKS, where: dict = None) -> list[Document]:
    identifiers = _extract_code_tokens(query)
    if not identifiers:
        return []

    corpus = _keyword_corpus(where=where)
    scored: list[tuple[Document, float]] = []
    for row in corpus["rows"]:
        haystack = row["haystack"]
        matched = [identifier for identifier in identifiers if _contains_exact_identifier(haystack, identifier)]
        if not matched:
            continue
        score = float(len(matched) * 100)
        first_positions = [haystack.find(identifier) for identifier in matched if haystack.find(identifier) >= 0]
        if first_positions:
            score += 1.0 / (min(first_positions) + 1)
        doc = Document(
            page_content=row["doc"].page_content,
            metadata={**row["doc"].metadata, "exact_code_match": ", ".join(matched)},
        )
        scored.append((doc, score))

    scored.sort(key=lambda item: item[1], reverse=True)
    return [doc for doc, _ in scored[:k]]


def _keyword_search_with_score(query: str, k: int = KEYWORD_POOL_SIZE, where: dict = None) -> list[tuple[Document, float]]:
    query_tokens = _tokenize(query)
    corpus = _keyword_corpus(where=where)
    rows = corpus["rows"]
    if not query_tokens or not rows:
        return []
    doc_freq = corpus["doc_freq"]
    total_docs = corpus["total_docs"]
    avg_doc_len = corpus["avg_doc_len"]
    idf = {
        token: math.log(1 + (total_docs - doc_freq.get(token, 0) + 0.5) / (doc_freq.get(token, 0) + 0.5))
        for token in set(query_tokens)
    }
    scored = []
    for row in rows:
        score = _keyword_score(query_tokens, row, idf, avg_doc_len)
        if score > 0:
            scored.append((row["doc"], score))
    scored.sort(key=lambda item: item[1], reverse=True)
    return scored[:k]


def _doc_key(doc: Document) -> tuple:
    return (
        os.path.basename(doc.metadata.get("source", "")),
        doc.metadata.get("chunk_index", doc.metadata.get("chunk_id", doc.page_content[:80])),
    )


def _hybrid_search(query: str, k: int = NUM_CHUNKS, where: dict = None) -> list[Document]:
    """
    Combine semantic vector retrieval with exact lexical matching.
    This improves part numbers, error codes, revision IDs, and table/spec lookups.
    """
    exact_hits = _exact_identifier_search(query, k=k, where=where)
    vector_hits = _similarity_search_with_score(query, k=max(k * VECTOR_POOL_FACTOR, k), where=where)
    keyword_hits = _keyword_search_with_score(query, k=max(KEYWORD_POOL_SIZE, k), where=where)
    if not vector_hits:
        fallback = [doc for doc, _ in keyword_hits]
        return _dedup_chunks(exact_hits + fallback)[:k]

    by_key: dict[tuple, dict] = {}
    max_keyword = max([score for _, score in keyword_hits], default=1.0)

    for rank, (doc, distance) in enumerate(vector_hits):
        key = _doc_key(doc)
        vector_score = 1.0 / (1.0 + max(distance, 0.0))
        by_key[key] = {
            "doc": doc,
            "vector_score": vector_score,
            "keyword_score": 0.0,
            "vector_rank_bonus": 1.0 / (rank + 1),
        }

    for doc, score in keyword_hits:
        key = _doc_key(doc)
        row = by_key.setdefault(key, {
            "doc": doc,
            "vector_score": 0.0,
            "keyword_score": 0.0,
            "vector_rank_bonus": 0.0,
        })
        row["keyword_score"] = max(row["keyword_score"], score / max_keyword)

    ranked = []
    for row in by_key.values():
        final_score = (
            HYBRID_ALPHA * row["vector_score"]
            + (1.0 - HYBRID_ALPHA) * row["keyword_score"]
            + 0.03 * row["vector_rank_bonus"]
        )
        row["doc"].metadata["retrieval_score"] = round(final_score, 6)
        row["doc"].metadata["keyword_score"] = round(row["keyword_score"], 6)
        ranked.append((row["doc"], final_score))

    ranked.sort(key=lambda item: item[1], reverse=True)
    hybrid_hits = [doc for doc, _ in ranked]
    return _dedup_chunks(exact_hits + hybrid_hits)[:k]


def _similarity_search(query: str, k: int = NUM_CHUNKS, where: dict = None) -> list[Document]:
    """Query ChromaDB and return LangChain Document objects."""
    return _hybrid_search(query, k=k, where=where)


def _vector_search(query: str, k: int = NUM_CHUNKS, where: dict = None) -> list[Document]:
    """Query ChromaDB and return LangChain Document objects."""
    n = min(k, collection.count())
    if n == 0:
        return []
    kwargs: dict = {"query_texts": [query], "n_results": n, "include": ["documents", "metadatas"]}
    if where:
        kwargs["where"] = where
    results = collection.query(**kwargs)
    return [Document(page_content=t, metadata=m)
            for t, m in zip(results["documents"][0], results["metadatas"][0])]


def _similarity_search_with_score(query: str, k: int = NUM_CHUNKS, where: dict = None) -> list[tuple[Document, float]]:
    """Query ChromaDB and return (Document, distance) tuples."""
    n = min(k, collection.count())
    if n == 0:
        return []
    kwargs: dict = {
        "query_texts": [query],
        "n_results": n,
        "include": ["documents", "metadatas", "distances"],
    }
    if where:
        kwargs["where"] = where
    results = collection.query(**kwargs)
    return [(Document(page_content=t, metadata=m), d)
            for t, m, d in zip(results["documents"][0], results["metadatas"][0], results["distances"][0])]


# ── Scope detection ──────────────────────────────────────────────────────────

_SCOPE_CANDIDATE_K      = 12
_SCORE_COMPETITION_GAP  = float(os.getenv("RAG_SCORE_COMPETITION_GAP", "0.15"))
_DEFAULT_MIN_RELEVANCE  = "0.55" if EMBEDDING_MODEL == "text-embedding-3-small" else "0.20"
_MIN_RELEVANCE_SCORE    = float(os.getenv("RAG_MIN_RELEVANCE_SCORE", _DEFAULT_MIN_RELEVANCE))
_DOMINANCE_RATIO        = float(os.getenv("RAG_DOMINANCE_RATIO", "0.10"))

# Words that occur in many filenames and carry no distinguishing meaning.
# Used by _extract_distinctive_keywords to avoid matching on 'manual', 'procedure', etc.
_FILENAME_STOP_WORDS = {
    "sop", "manual", "operation", "operations", "operator", "procedure",
    "schedule", "log", "inspection", "report", "document", "documents",
    "guide", "guidelines", "instruction", "instructions", "basic", "ref",
    "reference", "overview", "pdf", "doc", "txt", "md", "log", "csv", "tsv",
    "html", "htm", "json", "xml", "docx", "rtf", "xlsx", "xlsm", "pptx",
}


def _extract_distinctive_keywords(filename: str) -> set:
    """Pull meaningful words from a filename, dropping generic filler.

    Example: 'Forklift_Operation_Manual.pdf' → {'forklift'}
    """
    stem  = Path(filename).stem.lower()
    words = set(stem.replace("-", " ").replace("_", " ").split())
    # drop pure numbers (e.g. '001', '002', '2025')
    words = {w for w in words if not w.isdigit()}
    return words - _FILENAME_STOP_WORDS


def _display_name(filename: str) -> str:
    """'SOP-001-Assembly-Line-Startup.txt' → 'SOP-001 Assembly Line Startup'"""
    return Path(filename).stem.replace("-", " ").replace("_", " ")


def _source_matches_query(filename: str, query: str) -> bool:
    query_tokens = set(_tokenize(query))
    if not query_tokens:
        return False
    display = _display_name(filename).lower()
    display = re.sub(r"\bfound testing documentation\b", " ", display)
    display_tokens = [
        t for t in _tokenize(display)
        if t not in _FILENAME_STOP_WORDS and t not in {"found", "testing", "documentation"}
    ]
    if len(display_tokens) >= 2 and all(t in query_tokens for t in display_tokens[: min(4, len(display_tokens))]):
        return True
    distinctive = _extract_distinctive_keywords(filename) - {"found", "testing", "documentation"}
    overlap = {
        token for token in distinctive
        if any(q == token or q.startswith(token) or token.startswith(q) for q in query_tokens)
    }
    if len(overlap) >= 2:
        return True
    if len(overlap) == 1 and any(marker in query_tokens for marker in {"fda", "boeing", "autodesk", "modern"}):
        return True
    compact_display = " ".join(display_tokens)
    compact_query = " ".join(_tokenize(query))
    return bool(compact_display and compact_display in compact_query)


def _named_source_from_query(query: str) -> str | None:
    registry = _load_doc_registry()
    matches = [filename for filename in registry if _source_matches_query(filename, query)]
    if len(matches) == 1:
        return matches[0]
    if len(matches) > 1:
        query_tokens = set(_tokenize(query))
        ranked = sorted(
            matches,
            key=lambda filename: len((_extract_distinctive_keywords(filename) - {"found", "testing", "documentation"}) & query_tokens),
            reverse=True,
        )
        top = ranked[0]
        top_score = len((_extract_distinctive_keywords(top) - {"found", "testing", "documentation"}) & query_tokens)
        runner_score = len((_extract_distinctive_keywords(ranked[1]) - {"found", "testing", "documentation"}) & query_tokens)
        if top_score > runner_score:
            return top
    return None


def _build_clarification_options(sources: list) -> list:
    options = [{"label": _display_name(s), "value": s} for s in sources]
    options.append({"label": "All relevant documents", "value": "__all__"})
    return options


def _build_clarification_question(sources: list) -> str:
    names = [f'"{_display_name(s)}"' for s in sources[:3]]
    more  = f" (and {len(sources) - 3} more)" if len(sources) > 3 else ""
    if len(names) == 1:
        doc_list = names[0]
    elif len(names) == 2:
        doc_list = f"{names[0]} and {names[1]}"
    else:
        doc_list = f"{names[0]}, {names[1]}, and {names[2]}"
    return (
        f"I found relevant content in multiple documents: {doc_list}{more}. "
        f"Which one are you asking about, or would you like me to check all relevant documents?"
    )


def _detect_scope(message: str, session_id: str) -> tuple:
    """
    Determine retrieval scope. Returns (scope, data):
      ("pass",            None)             — proceed with full chain as-is
      ("broad",           None)             — explicit multi-doc synthesis
      ("ambiguous",       [filenames])      — multiple docs match, needs clarification
      ("resolved_single", (file, orig_q))  — clarification resolved to one doc
      ("resolved_all",    orig_q)          — clarification resolved to all docs
    """
    msg_lower = message.lower().strip()
    session   = conversation_sessions.get(session_id, {})
    pending   = session.get("pending_clarification")

    # 1 — Resolve a pending clarification
    if pending:
        choice = message.strip().lower()
        for opt in pending["options"]:
            label = str(opt.get("label", "")).strip().lower()
            value = str(opt.get("value", "")).strip().lower()
            display_value = _display_name(str(opt.get("value", ""))).lower()
            if choice in {label, value, display_value}:
                session["pending_clarification"] = None
                orig = pending["original_question"]
                if opt["value"] == "__all__":
                    return ("resolved_all", orig)
                return ("resolved_single", (opt["value"], orig))
        # User asked something new — clear pending and fall through
        session["pending_clarification"] = None

    # 2 — Explicit broad intent
    if any(phrase in msg_lower for phrase in BROAD_INTENT_PHRASES):
        return ("broad", None)

    # 3 — Candidate retrieval + score-based competition check
    named_source = _named_source_from_query(message)
    if named_source:
        return ("resolved_single", (named_source, message))

    candidates = _similarity_search_with_score(message, k=_SCOPE_CANDIDATE_K)
    if not candidates:
        return ("pass", None)

    # Track best (lowest L2) score and chunk list per source
    best_score:  dict = {}
    by_source:   dict = {}
    for doc, score in candidates:
        src = os.path.basename(doc.metadata.get("source", "unknown"))
        by_source.setdefault(src, []).append(doc)
        if src not in best_score or score < best_score[src]:
            best_score[src] = score

    if len(by_source) <= 1:
        return ("pass", None)

    top_score = min(best_score.values())

    # If user explicitly named a document (full stem), skip clarification
    for src in by_source:
        stem = Path(src).stem.lower()
        if stem in msg_lower or stem.replace("-", " ") in msg_lower or stem.replace("_", " ") in msg_lower:
            return ("pass", None)

    # Partial keyword match — if a distinctive word from exactly one filename
    # appears in the query AND that document is the best (or near-best) scorer
    # AND that keyword doesn't appear in the retrieved chunks of other sources
    # (which would mean it's a generic cross-cutting term, not a distinctive
    # entity name like "forklift" or "HPM100").
    msg_words = set(msg_lower.split())
    keyword_matches: dict = {}          # src → set of matched keywords
    for src in by_source:
        keywords = _extract_distinctive_keywords(src)
        # Match if a keyword is an exact word OR a prefix of a word in the query
        # (e.g. keyword "forklift" matches query word "forklifts")
        overlap = {
            kw for kw in keywords
            if any(w == kw or w.startswith(kw) or kw.startswith(w)
                   for w in msg_words)
        }
        if overlap:
            keyword_matches[src] = overlap
    if len(keyword_matches) == 1:
        matched_src   = next(iter(keyword_matches))
        matched_score = best_score[matched_src]
        matched_kw    = keyword_matches[matched_src]
        # Verify the keyword is truly distinctive: if it appears far more in the
        # matched doc's chunks than in all other docs' chunks combined, it's a
        # specific entity (like "forklift") — not a generic cross-cutting term.
        matched_text = " ".join(doc.page_content.lower() for doc in by_source[matched_src])
        other_text   = " ".join(
            doc.page_content.lower()
            for src2, docs in by_source.items() if src2 != matched_src
            for doc in docs
        )
        for kw in matched_kw:
            matched_count = matched_text.count(kw)
            other_count   = other_text.count(kw)
            # Distinctive if matched doc has the keyword and others have ≤25% as many occurrences
            if matched_count > 0 and other_count <= matched_count * 0.25:
                if matched_score <= top_score + 0.05:
                    return ("pass", None)
                break

    # If even the best match is a poor semantic fit, skip clarification — no doc
    # truly contains the answer so disambiguation wouldn't help.
    if top_score >= _MIN_RELEVANCE_SCORE:
        return ("pass", None)
    # Dominance check — if the top doc's score is substantially better than the
    # runner-up (by ratio), one doc clearly owns this query.
    sorted_scores = sorted(best_score.values())
    if len(sorted_scores) >= 2 and top_score > 0:
        runner_up = sorted_scores[1]
        if (runner_up - top_score) / top_score >= _DOMINANCE_RATIO:
            return ("pass", None)
    competitive  = [s for s, sc in best_score.items() if sc <= top_score + _SCORE_COMPETITION_GAP]

    if len(competitive) >= 2:
        ranked = sorted(competitive, key=lambda s: best_score[s])
        return ("ambiguous", ranked)

    return ("pass", None)


def _evict_stale_sessions() -> None:
    """Remove sessions that have been idle longer than SESSION_TTL_SECONDS."""
    now     = time.monotonic()
    stale   = [sid for sid, s in conversation_sessions.items()
               if now - s.get("last_accessed", now) > SESSION_TTL_SECONDS]
    for sid in stale:
        conversation_sessions.pop(sid, None)


def _answer_single_doc(question: str, filename: str, session_id: str):
    """Retrieve chunks from a specific file using ChromaDB metadata filter."""
    doc_chunks = _similarity_search(
        question,
        k=NUM_CHUNKS,
        where={"source": str(KNOWLEDGE_BASE_DIR / filename)},
    )

    if not doc_chunks:
        return jsonify({
            "reply":      f'The document "{_display_name(filename)}" does not appear to contain information about this.',
            "session_id": session_id,
            "metadata":   {"sources": []},
        })

    context          = _build_numbered_context(doc_chunks)
    result           = llm.invoke(QA_PROMPT.format(context=context, question=question))
    answer, used_chunks = _parse_citations(result.content, doc_chunks)

    sources = []
    for i, doc in enumerate(_dedup_chunks(used_chunks), 1):
        src     = doc.metadata.get("source", "Unknown")
        page    = doc.metadata.get("page", 0)
        snippet = doc.page_content[:150].replace("\n", " ")
        if len(doc.page_content) > 150:
            snippet += "..."
        sources.append({
            "id":         i,
            "file":       os.path.basename(src),
            "page":       page + 1,
            "section":    doc.metadata.get("section", ""),
            "doc_type":   doc.metadata.get("doc_type", ""),
            "revision":   doc.metadata.get("revision", ""),
            "status":     doc.metadata.get("status", ""),
            "snippet":    snippet,
            "full_snippet": doc.page_content.replace("\n", " "),
        })

    return jsonify({
        "reply":      answer,
        "session_id": session_id,
        "metadata":   {"sources": sources},
    })


def _answer_multi_doc(question: str, session_id: str):
    """Retrieve chunks balanced across sources and synthesize with multi-doc prompt."""
    # Scale k with the number of known documents so every doc gets fair representation.
    # Cap total chunks at MAX_MULTI_DOC_CHUNKS to stay within the LLM context window.
    num_docs       = max(1, len(_load_doc_registry()))
    pool_k         = min(num_docs * 5, 80)            # retrieve a wide pool
    chunks_per_src = max(2, MAX_MULTI_DOC_CHUNKS // num_docs)  # balance per source
    candidates     = _similarity_search(question, k=pool_k)
    by_source: dict = {}
    for doc in candidates:
        src = os.path.basename(doc.metadata.get("source", "unknown"))
        if src not in by_source:
            by_source[src] = []
        if len(by_source[src]) < chunks_per_src:
            by_source[src].append(doc)

    context_parts: list = []
    all_chunks:    list = []
    chunk_num = 1
    for src, chunks in by_source.items():
        context_parts.append(f"[Source: {src}]")
        for chunk in chunks:
            context_parts.append(f"[CHUNK {chunk_num}]\n{chunk.page_content}")
            all_chunks.append(chunk)
            chunk_num += 1
        context_parts.append("")

    result = llm.invoke(_MULTI_DOC_QA_TEMPLATE.format(
        context="\n".join(context_parts),
        question=question,
    ))
    answer, used_chunks = _parse_citations(result.content, all_chunks)

    sources = []
    for i, doc in enumerate(_dedup_chunks(used_chunks), 1):
        src     = doc.metadata.get("source", "Unknown")
        page    = doc.metadata.get("page", 0)
        snippet = doc.page_content[:150].replace("\n", " ")
        if len(doc.page_content) > 150:
            snippet += "..."
        sources.append({
            "id":         i,
            "file":       os.path.basename(src),
            "page":       page + 1,
            "section":    doc.metadata.get("section", ""),
            "doc_type":   doc.metadata.get("doc_type", ""),
            "revision":   doc.metadata.get("revision", ""),
            "status":     doc.metadata.get("status", ""),
            "snippet":    snippet,
            "full_snippet": doc.page_content.replace("\n", " "),
        })

    return jsonify({
        "reply":      answer,
        "session_id": session_id,
        "metadata":   {"sources": sources},
    })


# ── Helper: manual conversational RAG pipeline ───────────────────────────────
def _chat_with_memory(question: str, session_id: str) -> tuple[str, list[Document]]:
    """
    Manual replacement for ConversationalRetrievalChain.
    1. Condense follow-up questions using CONDENSE_PROMPT + chat history.
    2. Retrieve top-k chunks from ChromaDB.
    3. Answer with QA_PROMPT (strict grounding).
    4. Save to ConversationBufferMemory for next turn.
    """
    session = conversation_sessions.setdefault(session_id, {})
    session["last_accessed"] = time.monotonic()

    if "memory" not in session:
        session["memory"] = ConversationBufferMemory(
            memory_key="chat_history",
            return_messages=False,
        )

    memory      = session["memory"]
    chat_history = memory.load_memory_variables({}).get("chat_history", "")

    # Step 1: Condense follow-up question into standalone form if needed
    if chat_history:
        condensed = llm.invoke(
            CONDENSE_PROMPT.format(chat_history=chat_history, question=question)
        ).content.strip()
    else:
        condensed = question

    # Step 2: Retrieve relevant chunks
    chunks = _similarity_search(condensed, k=NUM_CHUNKS)
    if not chunks:
        answer = "The uploaded documents do not contain information about this."
        memory.save_context({"input": question}, {"output": answer})
        return answer, []

    # Step 3: Answer strictly from retrieved context
    context   = _build_numbered_context(chunks)
    raw       = llm.invoke(QA_PROMPT.format(context=context, question=condensed)).content
    answer, chunks = _parse_citations(raw, chunks)

    # Step 4: Persist turn to memory
    memory.save_context({"input": question}, {"output": answer})
    return answer, chunks


# ═══════════════════════════════════════════════════════════════════════════════
#  Routes
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/chat", methods=["POST"])
def chat():
    """
    Main chat endpoint — matches the contract the Next.js proxy expects.
    Body:   { message: str, history: list, session_id?: str }
    Reply:  { reply: str, session_id: str, metadata: { sources: [] } }
    """
    try:
        data       = request.get_json(silent=True) or {}
        message    = (data.get("message") or "").strip()
        session_id = data.get("session_id") or "default"

        if not message:
            return jsonify({"error": "Message cannot be empty"}), 400

        # Evict idle sessions to prevent unbounded memory growth
        _evict_stale_sessions()

        # Lazy-init ChromaDB collection
        if collection is None:
            ok, msg = _init_store()
            if not ok:
                return jsonify({
                    "reply":      "The assistant isn't ready yet. Please upload a document first.",
                    "session_id": session_id,
                    "metadata":   {"sources": []},
                }), 503

        # Guard: no documents uploaded yet
        if collection.count() == 0:
            return jsonify({
                "reply":      "No documents have been uploaded yet. Please upload a document using the sidebar before asking questions.",
                "session_id": session_id,
                "metadata":   {"sources": []},
            })

        # Guard: block personal/entertainment/small-talk questions BEFORE retrieval.
        # Still skip when a clarification is pending — the user is replying to a
        # document-selection prompt, so their short answer ("the first one", "yes")
        # should never be classified as off-topic.
        session_data = conversation_sessions.get(session_id, {})
        has_history = "memory" in session_data and session_data["memory"].load_memory_variables({}).get("chat_history", "")
        has_pending_clarification = bool(session_data.get("pending_clarification"))
        if not has_pending_clarification and _is_off_topic(message):
            return jsonify({
                "reply":      "I can only answer questions about the uploaded documents.",
                "session_id": session_id,
                "metadata":   {"sources": []},
            })

        # ── Condense follow-ups before scope detection ──────────────────────────
        # If there's conversation history, rephrase the message into a standalone
        # question so that scope detection works on the full context, not just
        # a bare pronoun like "tell me more about the first one."
        chat_history = ""
        if has_history:
            chat_history = session_data["memory"].load_memory_variables({}).get("chat_history", "")
            condensed_for_scope = llm.invoke(
                CONDENSE_PROMPT.format(chat_history=chat_history, question=message)
            ).content.strip()
        else:
            condensed_for_scope = message

        # ── Scope detection ──────────────────────────────────────────────────────
        scope, scope_data = _detect_scope(condensed_for_scope, session_id)

        if scope == "ambiguous":
            options  = _build_clarification_options(scope_data)
            question = _build_clarification_question(scope_data)
            conversation_sessions.setdefault(session_id, {})["pending_clarification"] = {
                "original_question": message,
                "options":           options,
            }
            return jsonify({
                "reply":      question,
                "session_id": session_id,
                "metadata":   {
                    "sources":       [],
                    "clarification": {"question": question, "options": options},
                },
            })

        if scope in ("broad", "resolved_all"):
            return _answer_multi_doc(scope_data if scope == "resolved_all" else message, session_id)

        if scope == "resolved_single":
            filename, original_q = scope_data
            return _answer_single_doc(original_q, filename, session_id)

        # ── Default: manual conversational RAG pipeline ─────────────────────────
        answer, source_docs = _chat_with_memory(message, session_id)

        sources = []
        for i, doc in enumerate(_dedup_chunks(source_docs), 1):
            src     = doc.metadata.get("source", "Unknown")
            page    = doc.metadata.get("page", 0)
            snippet = doc.page_content[:150].replace("\n", " ")
            if len(doc.page_content) > 150:
                snippet += "..."
            sources.append({
                "id":         i,
                "file":       os.path.basename(src),
                "page":       page + 1,
                "section":    doc.metadata.get("section", ""),
                "doc_type":   doc.metadata.get("doc_type", ""),
                "revision":   doc.metadata.get("revision", ""),
                "status":     doc.metadata.get("status", ""),
                "snippet":    snippet,
                "full_snippet": doc.page_content.replace("\n", " "),
            })

        return jsonify({
            "reply":      answer,
            "session_id": session_id,
            "metadata":   {"sources": sources},
        })

    except Exception as exc:
        traceback.print_exc()
        return jsonify({"error": str(exc)}), 500


@app.route("/api/documents", methods=["GET"])
def list_documents():
    """Return a list of all uploaded/ingested documents."""
    registry = _load_doc_registry()
    docs = [
        {
            "filename": name,
            "chunks": info.get("chunks", 0),
            "uploaded_at": info.get("uploaded_at", ""),
            "doc_type": info.get("doc_type", ""),
            "revision": info.get("revision", ""),
            "status": info.get("status", ""),
            "embedding_model": info.get("embedding_model", ""),
        }
        for name, info in registry.items()
    ]
    return jsonify({"documents": docs})


@app.route("/api/documents/upload", methods=["POST"])
def upload_document():
    """Upload and ingest a supported document file."""
    if "file" not in request.files:
        return jsonify({"success": False, "message": "No file field in request"}), 400

    file = request.files["file"]
    ext = Path(file.filename).suffix.lower() if file.filename else ""
    if not file.filename or ext not in ALLOWED_EXTENSIONS:
        allowed = ", ".join(sorted(ALLOWED_EXTENSIONS))
        return jsonify({"success": False, "message": f"Unsupported file type. Allowed: {allowed}"}), 400

    filename = secure_filename(Path(file.filename).name)
    filepath = KNOWLEDGE_BASE_DIR / filename
    file.save(filepath)

    # Ensure ChromaDB collection is initialised before ingestion
    if collection is None:
        ok, msg = _init_store()
        if not ok:
            filepath.unlink(missing_ok=True)
            return jsonify({"success": False, "message": msg}), 503

    ok, msg, chunk_count = _ingest_file(filepath)
    if not ok:
        filepath.unlink(missing_ok=True)
        return jsonify({"success": False, "message": msg}), 500

    # Update registry
    registry = _load_doc_registry()
    doc_meta = {}
    try:
        sample_docs = _load_docs(filepath)
        doc_meta = _derive_doc_metadata(filepath, sample_docs)
    except Exception:
        doc_meta = {"embedding_model": EMBEDDING_MODEL}
    registry[filename] = {
        "chunks": chunk_count,
        "uploaded_at": datetime.now(timezone.utc).isoformat(),
        "doc_type": doc_meta.get("doc_type", ""),
        "revision": doc_meta.get("revision", ""),
        "status": doc_meta.get("status", ""),
        "embedding_model": EMBEDDING_MODEL,
    }
    _save_doc_registry(registry)

    # Clear sessions so they don't use stale context
    conversation_sessions.clear()

    return jsonify({"success": True, "message": msg, "filename": filename, "chunks": chunk_count})


@app.route("/api/documents/<filename>", methods=["DELETE"])
def delete_document(filename: str):
    """
    Delete a document file and remove its vectors from ChromaDB.
    ChromaDB supports selective deletion by metadata — no full index rebuild needed.
    """
    registry = _load_doc_registry()
    if filename not in registry:
        return jsonify({"success": False, "message": "Document not found"}), 404

    # Remove the file from disk
    (KNOWLEDGE_BASE_DIR / filename).unlink(missing_ok=True)

    # Remove from registry
    del registry[filename]
    _save_doc_registry(registry)

    # Selectively delete only this document's vectors — no rebuild required
    collection.delete(where={"source": str(KNOWLEDGE_BASE_DIR / filename)})
    _keyword_corpus_cache.clear()

    # Clear sessions so memory doesn't reference deleted content
    conversation_sessions.clear()

    return jsonify({"success": True, "message": f"{filename} deleted successfully."})


@app.route("/api/clear", methods=["POST"])
def clear_conversation():
    """Clear conversation history for a session."""
    data       = request.get_json(silent=True) or {}
    session_id = data.get("session_id", "default")
    conversation_sessions.pop(session_id, None)
    return jsonify({"status": "success", "message": "Conversation cleared."})


@app.route("/api/health", methods=["GET"])
def health():
    """Health / readiness check."""
    has_api_key  = bool(os.getenv("OPENAI_API_KEY"))
    has_docs     = collection is not None and collection.count() > 0
    return jsonify({
        "status":          "healthy",
        "has_documents":   has_docs,
        "has_api_key":     has_api_key,
        "ready":           has_docs and has_api_key and collection is not None,
        "active_sessions": len(conversation_sessions),
        "embedding_model": EMBEDDING_MODEL,
        "collection":      COLLECTION_NAME,
    })


# ── Startup ──────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("\n" + "=" * 60)
    print("  Kuldeep RAG Chatbot  —  Flask Backend (ChromaDB)")
    print("=" * 60)
    _init_store()   # best-effort on startup; 
    print("  Open: http://localhost:5000")
    print("=" * 60 + "\n")
    app.run(debug=True, host="0.0.0.0", port=5000)

